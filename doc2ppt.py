import pptx as ppt
from pptx.util import Cm
import docx as word
import re

lineSpace = float(input('行间距: '))
fontsize = int(input('字体大小 (0将不会改变): '))

# read content from docx files
filename = input('文档名称: ')
rDocSoc = word.Document(filename)
wPPTSoc = ppt.Presentation()
# create slide
slideLayout = wPPTSoc.slide_layouts.get_by_name('Blank')
slide = wPPTSoc.slides.add_slide(slideLayout)
left = Cm(1.94)
top = Cm(1.15)
width = Cm(20.52)
height = Cm(10)
txBox = slide.shapes.add_textbox(left, top, width, height)
tf = txBox.text_frame

PPTCount = 0

def getStyle(p, check=False):
    for run in p.runs:
        if check == True:
            if run.font.name == None or run.font.size == None:
                continue
            else:
                return run.font.name, run.font.size
        else:
            return run.font.name, run.font.size

for p in rDocSoc.paragraphs:
    if '++++++++++' in p.text:
        PPTCount += 1
        slideLayout = wPPTSoc.slide_layouts.get_by_name('Blank')
        slide = wPPTSoc.slides.add_slide(slideLayout)
        left = Cm(1.94)
        top = Cm(1.15)
        width = Cm(20.52)
        height = Cm(10)
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame

    else:
        # add text content
        slideP = tf.add_paragraph()
        if len(p.text) > 30:
            slideP.text = '\n'.join(re.findall(r'.{34}', p.text))
        else:
            slideP.text = p.text
        try:
            if fontsize == 0:
                slideP.font.font, slideP.font.size = getStyle(p, True)
            else:
                slideP.font.font, slideP.font.size = getStyle(p, True)
                slideP.font.size = fontsize
        except:
            slideP.font.font = getStyle(p)
        slideP.line_spacing = lineSpace

wPPTSoc.save('result.pptx')
